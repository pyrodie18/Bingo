from pptx import Presentation
from pptx.util import Inches
from random import randint
from random import seed
import time



def BuildBoard(table):
    seed()
    SeedVal = randint(1,100000)
    seed(SeedVal)

    table.cell(0,0).text = "B"
    table.cell(0,1).text = "I"
    table.cell(0,2).text = "N"
    table.cell(0,3).text = "G"
    table.cell(0,4).text = "O"
    FillNumbers(table, 0, 1, 15)
    FillNumbers(table, 1, 16, 30)
    FillNumbers(table, 2, 31, 45)
    FillNumbers(table, 3, 46, 60)
    FillNumbers(table, 4, 61, 75)
    table.cell(3,2).text = "FREE"
    table.cell(6,0).merge(table.cell(6,4))
    table.cell(6,0).text = str(SeedVal)

def FillNumbers(table, Col, Low, High):
    values = []
    for i in range (1,6):
        while True:
            value = randint(Low,High)
            if not(value in values):
                values.append(value)
                break
        table.cell(i,Col).text = str(value)



rows = 7
cols = 5
left = top = Inches(0.5)
width = Inches(4.0)
height = Inches(3.0)

for i in range (10):
    prs = Presentation()
    title_only_slide_layout = prs.slide_layouts[6]

    for slide in range (4):
        CurrentSlide = prs.slides.add_slide(title_only_slide_layout)
        shapes = CurrentSlide.shapes

        table1 = shapes.add_table(rows, cols, Inches(0.5), Inches(0.5), width, height).table
        table2 = shapes.add_table(rows, cols, Inches(5.5), Inches(0.5), width, height).table
        table3 = shapes.add_table(rows, cols, Inches(0.5), Inches(4.0), width, height).table
        table4 = shapes.add_table(rows, cols, Inches(5.5), Inches(4.0), width, height).table

        BuildBoard(table1)
        BuildBoard(table2)
        BuildBoard(table3)
        BuildBoard(table4)

    filename = "board" + str(i) + ".pptx"
    prs.save(filename)
